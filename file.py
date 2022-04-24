import directories

from pdfminer.high_level import extract_text
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
import docx
import pptx
from bs4 import BeautifulSoup
from urllib.request import Request, urlopen
from urllib.parse import urlparse

import os
import re
import pickle

import sqlite3



class File_base(object):
    """docstring for File_base."""

    def __init__(self, path, contents_type):
        # super(File_base, self).__init__(path)
        self.path = path
        # empty mean text and b mean binary(python)
        self.contents_type = contents_type
        assert self.is_file(), f"path '{path}' is not a valid file path"

    def read(self, mode):
        file_obj = self.get_file_object(mode)
        contents = file_obj.read()
        file_obj.close()
        return contents

    def is_file_empty(self):
        return directories.is_file_empty(self.path)

    def get_file_object(self, mode):
        return open(self.path, mode,  encoding='utf8')

    def is_file(self):
        return directories.is_file(self.path)

    def get_filename_ext(self):
        '''Generates filename for text file'''
        filename_ext = directories.get_filename_extension(self.path)
        return filename_ext

    def get_filename(self):
        filename = directories.get_filename(self.path)
        return filename

    def get_path(self):
        return self.path

    def set_path(self, path):
        self.path = path

    def files_paths(folder_path, format=None):
        return directories.files_paths(folder_path, format)

    def clear_filename(self, filename):
        '''converts invalid string into sring that can be used for filenames.
        Flask has secure_filename() function but dependency on flask will cause problems.'''
        filename = directories.get_valid_filename(filename)
        if len(filename) >= 100:
            return filename[:99]
        else:
            return filename



class DB_file(File_base):
    """docstring for DB_file."""

    def __init__(self, path=None, check_same_thread=True):
        if path != None and path != ":memory:":
            super(DB_file, self).__init__(path, "b")
        else:
            self.path = ":memory:"
        self.check_same_thread = check_same_thread
        self.conn = self.create_conn(self.path, check_same_thread)
        self.commit()

    def get_conn(self):
        return self.conn

    def create_conn(self, path, check_same_thread):
        #beaware of last argumant(threads safty issues if data not managed)
        #conn = sqlite3.connect(path, check_same_thread)
        conn = sqlite3.connect(path, timeout=40)
        conn.row_factory = sqlite3.Row
        return conn

    def get_cursor(self):
        return self.conn.cursor()

    def table_exists(self, table_name):
        cursor = self.get_cursor()
        cursor.execute(f"SELECT count(name) FROM sqlite_master WHERE type='table' AND name='{table_name}'")
        return cursor.fetchone()[0]==1

    def execute_command(self, command):
        with sqlite3.connect(self.path) as con:
            cur = con.cursor()
            cur.execute(command)
            con.commit()

    def to_db_compatible(self, collection, data_type=""):
        string_version = f"{tuple(collection)}"
        compatible = string_version.replace(',)', ')')
        if data_type:
            compatible = compatible.replace(",", f" {data_type},")
            compatible = compatible.replace("'", "").replace(")", f" {data_type})")
        return compatible

    def commit(self):
        self.conn.commit()

    def close(self):
        self.conn.close()

    def unlock_db(self):
        """Replace db_filename with the name of the SQLite database."""
        self.conn = self.create_conn()
        self.conn.commit()
        self.conn.close()


class DB_table(DB_file):
    """docstring for DB_table."""

    def __init__(self, conn, table_name, columns, primary_column):
        #super(DB_table, self).__init__(path)
        self.conn = conn
        self.table_name = table_name
        self.columns = columns
        self.primary_column = primary_column
        self.table_create_query = None
        #assert self.table_exists(table_name), f"'{table_name}' does not exist"
        self.create_table()

    def create_table(self):
        assert False, 'create_table method not implemented'

    def get_specific(self, column, items):
        compatible = self.to_db_compatible(items)
        cursor = self.get_cursor()
        query = f'''SELECT * FROM {self.table_name} WHERE {column} IN {compatible};'''
        cursor.execute(query)
        return self.get_dict_rows(cursor.fetchall())

    def delete_specific(self, column, items):
        compatible = self.to_db_compatible(items)
        cursor = self.get_cursor()
        query = f'''DELETE FROM {self.table_name} WHERE {column} IN {compatible};'''
        cursor.execute(query)
        self.conn.commit()

    def delete_all(self):
        cursor = self.get_cursor()
        query = f'''DELETE FROM {self.table_name}'''
        cursor.execute(query)
        self.conn.commit()

    def ratio_delete(self, ratio):
        assert ratio <= 1 and ratio >= 0, f"ration is out of range {ratio}"
        primary_items = self.get_primary_items(self)
        delete_primary_items = primary_items[0:int(len(primary_items)*ratio)]
        self.delete_specific(self, column, items)

    def get_column_items(self, column):
        cursor = self.get_cursor()
        query = f'''SELECT {column}  FROM {self.table_name};'''
        cursor.execute(query)
        dict_rows = self.get_dict_rows(cursor.fetchall())
        return self.get_dict_rows_values(dict_rows)

    def get_primary_items(self):
        return self.get_column_items(self.primary_column)

    def db_insert(self, dictionary):
        assert False, "db_inset() not implemented"

    def db_update(self, dictionary):
        assert False, "db_update not implemented"

    def db_clear_table(self):
        self.conn.execute(f"DELETE from {self.table_name};")

    def get_rows(self):
        cursor = self.conn.cursor()
        cursor.execute(f"select * from {self.table_name}")
        rows = cursor.fetchall();
        return rows

    def get_dict_rows(self, rows=None):
        if rows==None:
            rows = self.get_rows()
        return [dict(row) for row in rows]

    def concat_dict_rows(self, dict_rows):
        return {key : value for row in dict_rows for key, value in row.items()}

    def get_dict_rows_keys(self, dict_rows):
        return [key for row in dict_rows for key in row.keys()]

    def get_dict_rows_values(self, dict_rows):
        return [values for row in dict_rows for values in row.values()]

    def get_length(self):
        return len(self.get_column_items(self.primary_column))



class Map_data(object):
    """docstring for Map_data."""

    def __init__(self, data):
        super(Map_data, self).__init__()
        self.data = data

    def overide_map(self, map_like):
        self.data.clear()
        self.concat(map_like)

    def get_map(self):
        return self.data.copy()

    def is_map_empty(self):
        return len(self.data) == 0

    def get_specific_items(self, keys):
        '''Method tested'''
        return {key: self.data[key] for key in self.data.keys()
                               and keys if key in self.data}

    def get_items(self):
        return self.data.items()

    def get_values(self, key):
        return self.data.values()

    def get_keys(self):
        return self.data.keys()

    def get_value(self, key):
        return self.data[key]

    def is_key_avail(self, key):
        if key in self.data:
            return True
        else:
            return False

    def get_key(self, value_arg):
        for key, value in self.get_items():
            if value == value_arg:
                return key

    def append(self, key, value):
        self.data[key] = value

    def concat(self, map_like):
        for key, value in map_like.items():
            self.append(key, value)

    def remove(self, key):
        self.data.pop(key)



class Reload_table(DB_table, Map_data):
    """docstring for Reload_table."""

    def __init__(self, conn, table_name, columns, primary_column):
        """docstring for Extract_json."""
        super(Reload_table, self).__init__(conn, table_name, columns, primary_column)
        self.data = dict()
        #self.primary_column = {}


    def db_sync(self, primary_column_elements):
        database_primary = self.get_in_db_only(primary_column_elements)
        db_data = self.get_specific(self.primary_column, database_primary)
        for data in db_data:
            assert data[self.primary_column] not in self.data, f"{data[self.primary_column]} already in memory"
            self.concat(self.load(data))
        assert len(primary_column_elements) >= len(db_data), f'''updated db data \ {len(db_data)} is greater than expected >={len(primary_column_elements)}'''
        return list(database_primary)

    def sync_all(self):
        for load_data in self.load_all():
            self.concat(load_data)

    def get_in_db_only(self, elements, column=None):
        if column == None:
            column = self.primary_column
        db_primary_keys =  set(self.get_column_items(column))
        return [element for element in elements if not self.is_key_avail(element) and element in db_primary_keys]



    def load(self, db_item):
        #db_item --> sqlite.RAW object
        '''load contents from the database table converting them to python data strutures(objects) and then concatinating them into current data
        '''
        assert False, "load method not implemented"


    def dump(self, data_item):
        '''Convert from python object into format compatible with database table. The output is dictionaries within list that can be saved using db.db_insert(). It should also filter the results as py may contain items already in db.
        '''
        assert False, "dump method not implemented"

    def dumps(self, data):
        return [self.dump(data_item) for data_item in data]

    def loads(self, table_data):
        return [self.load(table_item) for table_item in table_data]

    def load_all(self):
        #get all items in db(you can optimize(no time))
        return self.loads(self.get_dict_rows())

    def filter_save_dumps(self, data):
        filtered = [data_item for data_item in data if self.is_save_ready(data_item)]
        dumped_filter = self.dumps(filtered)
        return dumped_filter

    def clear(self):
        self.data.clear()

    def is_db_synced(self):
        if len(self.data) == 0:
            return True
        primary_db_items = set(self.get_column_items(self.primary_column))
        current_keys = set(self.data.keys())
        return primary_db_items.issuperset(current_keys)


    def is_save_ready(self, source):
        #data_item represent an item within self.data
        '''body provided by subclass'''
        column_items = self.get_column_items(self.primary_column)
        return source not in column_items
        #assert False, "is_save_ready method not implemented"



    def save(self, overide=False):
        '''saves data into database'''
        if len(self.data) == 0:
            return

        if overide:
            dumped_list = self.dumps(self.data)
        else:
            dumped_list = self.filter_save_dumps(self.data)
        for element in dumped_list:
            if overide:
                self.db_update(element)
            else:
                self.db_insert(element)
        self.commit()
        assert self.is_db_synced(), "Not synced with database table after commiting."


class Extract_text(object):
    """docstring for ."""

    def __init__(self, source):
        super(Extract_text, self).__init__()
        self.source = source
        self.text = ''

    def is_valid(self):
        '''implemented by subclass'''

    def is_converted(self):
        if self.text == '':
            return False
        else:
            return True

    def convert(self):
        '''Code for converting file in the path to text'''
        pass

    def remove_multi_newlines(self, text):
        return re.sub(r'(\n+)+', r'\n', text)

    def remove_html(self, text):
        return BeautifulSoup(text, "html.parser").get_text()

    def process_text(self, text):
        html_removed = self.remove_html(text)
        multi_newlines_removed = self.remove_multi_newlines(html_removed)
        return multi_newlines_removed

    def save(self, folder_path):
        if not self.is_converted() and self.is_valid():
            self.convert()
        filename = self.clear_filename(self.get_filename_ext())
        file_path = os.path.join(folder_path, filename + ".txt")
        directories.write_file(file_path, self.text)
        return file_path

    def get_text(self):
        if not self.is_converted() and self.is_valid():
            self.convert()
        return self.process_text(self.text)


class File_to_text(File_base, Extract_text):
    """docstring for File_to_text."""

    def __init__(self, path, contents_type):
        super(File_to_text, self).__init__(path, contents_type)
        Extract_text.__init__(self, path)
        self.path = path

    def is_valid(self):
        '''validate the file in the path and return boolean'''
        return True


class Text_file_to_text(File_to_text):
    """docstring for Pdf_to_text."""

    def __init__(self, path, contents_type=""):
        super(Text_file_to_text, self).__init__(path, contents_type)

    def is_valid(self):
        try:
            #performs full read but reading text file is not competer intensive
            self.read("r"+self.contents_type)
            return True
        except:
            return False

    def convert(self):
        self.text = self.read("r"+self.contents_type)
        return self.text




class Pdf_to_text(File_to_text):
    """docstring for Pdf_to_text."""

    def __init__(self, path, contents_type="b"):
        super(Pdf_to_text, self).__init__(path, contents_type)

    def is_valid(self):
        with open(self.path, "rb") as file:
            try:
                parser = PDFParser(file)
                PDFDocument(parser)
                return True
            except:
                return False

    def convert(self):
        self.text = extract_text(self.path)
        return self.text


class Docx_to_text(File_to_text):
    """docstring for Docx_to_text."""

    def __init__(self, path, contents_type="b"):
        super(Docx_to_text, self).__init__(path, contents_type)

    def is_valid(self):
        try:
            docx.Document(self.path)
            return True
        except:
            return False

    def convert(self):
        text = ''
        document = docx.Document(self.path)
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        self.text = text
        return self.text


class Pptx_to_text(File_to_text):
    """docstring for Docx_to_text."""

    def __init__(self, path, contents_type="b"):
        super(Pptx_to_text, self).__init__(path, contents_type)

    def is_valid(self):
        try:
            pptx.Presentation(self.path)
            return True
        except:
            return False

    def convert(self):
        text = ""
        prs = pptx.Presentation(self.path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        self.text = text
        return self.text

class Url_to_text(Extract_text):
    """docstring for Docx_to_text."""

    def __init__(self, url, contents_type=""):
        self.url = url
        self.request = Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        super(Url_to_text, self).__init__(url)

    def is_valid(self):
        if urlparse(self.url).netloc == "datanavy.site":
            #crowling the application url restarts the application
            #Caused maybe by user-agent technologies
            return False
        req = Request(self.url, headers={'User-Agent': 'Mozilla/5.0'})
        try:
            # you can perfor validation based on returned code mostly 200
            code = urlopen(req).getcode()
            if code == 200:
                return True
            else:
                return False
        # except URLError
        # except urllib.request.URLError:
        except:
            return False

    def convert(self):
        assert self.is_valid(), f'The URL({self.source}) seems to be invalid'
        req = Request(self.url, headers={'User-Agent': 'Mozilla/5.0'})
        #url error here
        page = urlopen(req)
        #UnicodeDecodeError
        try:
            html = page.read().decode("utf-8")
        except UnicodeDecodeError as e:
            return None
        soup = BeautifulSoup(html, "html.parser")
        body = soup.find('body')
        text = re.sub(r'(\n+)+', r'\n', body.get_text())
        self.text = text
        return self.text

    def get_filename(self):
        return self.source

    def get_filename_ext(self):
        return self.source


class Recontructor(object):
    """docstring for Recontructor."""

    def __init__(self, data=None):
        super(Recontructor, self).__init__()
        self.data = data

    def get_File_type(self, filename):
        supported = [".pdf",".pptx",".docx"]
        extension = directories.get_file_extension(filename)
        assert self.is_filename(filename), f"seems as filename({filename}) is\
        not a valid filename"
        if extension in supported:
            return extension.replace(".","")
        else:
            return "text"


    def is_filename(self, source):
        extension = directories.get_file_extension(source)
        return not self.is_url(source) and bool(extension)

    def is_url(self, source):
        url_parse = urlparse(source)
        return bool(url_parse.netloc and url_parse.scheme)

    def is_text(self, source):
        return not (self.is_filename(source) or self.is_url(source))

    def get_source_type(self, source):
        if self.is_filename(source):
            return "file"
        elif self.is_url(source):
            return "webpage"
        elif self.is_text(source):
            return "text"
        assert False, f"cannot determine type of source '{source}'"




class Knowledge_text_map(Reload_table, Recontructor):
    """Basic maping of database table. Does not require special conversion of data"""

    def __init__(self, conn, table_name="knowledge_text", columns=["source_type", "text"], primary_column= "source"):
        super(Knowledge_text_map, self).__init__(conn, table_name, columns, primary_column)


    def create_table(self):
        if not self.table_exists(self.table_name):
            query = f'''CREATE TABLE {self.table_name}
             (source TEXT PRIMARY KEY     NOT NULL,
             text           TEXT    NOT NULL,
             source_type     TEXT);'''
            self.conn.execute(query)

    def db_insert(self, element):
        query = f'''INSERT INTO {self.table_name} (source, text)
      VALUES (?, ?)'''
        self.conn.execute(query, (element["source"], element["text"]))

    def load(self, db_item):
        return {db_item["source"]: db_item["text"]}


    def dump(self, data_item):
        assert data_item, f"seems as data item is empty its value is '{data_item}'"
        source = data_item
        text = self.data[source]
        return {"source": source, "text": text}

    def get_typed_data(self, dict_rows=None):
        '''Add knowledge type using prediction'''
        if dict_rows == None:
            dict_rows = self.get_dict_rows()
        copied_rows = dict_rows.copy()
        for dict_row in copied_rows:
            dict_row["source_type"] = self.get_source_type(dict_row['source'])
            if dict_row["source_type"] == "file":
                dict_row["file_type"] = self.get_File_type(dict_row['source'])
        return copied_rows


class Knowledge_bytes_map(Reload_table):
    """docstring for Bytes_file. For working with python bytes objects using pickle module"""

    def __init__(self, conn, table_name="knowledge_bytes", columns=["knowledge_bytes"], primary_column= "source"):
        super(Knowledge_bytes_map, self).__init__(conn, table_name, columns, primary_column)

    def create_table(self):
        if not self.table_exists(self.table_name):
            query = f'''CREATE TABLE {self.table_name}
             (source TEXT PRIMARY KEY     NOT NULL,
             knowledge_bytes           BLOB    NOT NULL,
             source_type     TEXT);'''
            self.conn.execute(query)

    def db_insert(self, element):
        query = f'''INSERT INTO {self.table_name} (source, knowledge_bytes)
      VALUES (?, ?)'''
        self.conn.execute(query, (element["source"], element["knowledge_bytes"]))

    def load(self, db_item):
        loaded_data = pickle.loads(db_item["knowledge_bytes"])
        return {db_item["source"]: loaded_data}

    def dump(self, data_item):
        assert data_item, f"seems as data item is empty its value is '{data_item}'"
        source = data_item
        knowledge_obj = self.data[source]
        dumped = pickle.dumps(knowledge_obj)
        return {"source": source, "knowledge_bytes": dumped}


class User_data(Reload_table):
    """Stores data for tracking a user into database"""

    def __init__(self, conn, table_name="user_data", columns=[ "attribute_value"], primary_column= "attribute_name"):
        super(User_data, self).__init__(conn, table_name, columns, primary_column)

    def create_table(self):
        if not self.table_exists(self.table_name):
            query = f'''CREATE TABLE {self.table_name}
             (attribute_name TEXT PRIMARY KEY     NOT NULL,
             attribute_value TEXT    NOT NULL
             );'''
            self.conn.execute(query)
            self.commit()

    def db_insert(self, dumped_element):
        query = f'''INSERT INTO {self.table_name} (attribute_name, attribute_value)
      VALUES (?, ?)'''
        self.conn.execute(query, (dumped_element["attribute_name"], dumped_element["attribute_value"]))

    def db_update(self, dumped_element):
        key = dumped_element['attribute_name']
        value = dumped_element['attribute_value']
        query = f"UPDATE {self.table_name} set attribute_value = '{value}' where attribute_name = '{key}';"
        self.conn.execute(query)

    def load(self, db_item):
        return {db_item["attribute_name"]: db_item["attribute_value"]}


    def dump(self, data_item):
        assert data_item, f"seems as data item is empty its value is '{data_item}'"
        attribute_name = data_item
        attribute_value = self.data[attribute_name]
        return {"attribute_name": attribute_name, "attribute_value": attribute_value}

def create_db(path):
    conn = sqlite3.connect(path, timeout=10)
    return conn

def unlock_db(db_filename):
    """Replace db_filename with the name of the SQLite database."""
    connection = sqlite3.connect(db_filename)
    connection.commit()
    connection.close()


if __name__ == "__main__":
    rec_obj = Recontructor(33)
    rec_obj.get_File_type("ppddd.ttrrr")


if __name__ == "__main__":
    path = r"test.db"
    #conn = create_db(path)
    db_obj = DB_file(path)
    conn = db_obj.get_conn()
    conn.commit()
    #conn.close()

    #conn.execute("INSERT INTO user_dat (attribute_name, attribute_value) \
    #  VALUES ('user_id', 'wewqq')")
    #db_obj = DB_file(path)
    #conn = db_obj.get_conn()
    # db_table = DB_table(db_obj.get_conn(), "__3meoeether_fucking", ["gender", "sport", "education"],
    #                     "id number")
    # db_table.to_db_compatible(["name", "village", "age", "education"], "TEXT")
    # db_table.get_column_items("gender")
    # db_table.db_insert({"gender": "ability sekgobela", "sport": "phlaborwa", "education": "matric madie"})
    # rows = db_table.get_rows()
    # dict_rows = db_table.get_dict_rows(rows)
    # db_table.get_dict_rows_values(dict_rows)

    Text_map =User_data(conn, "user_data")
    #Text_map.db_insert({'attribute_value':" guy is body", "attribute_name": "man is rich 233"})
    Text_map.db_update({'attribute_value':23333, "attribute_name": "man is rich 233"})
    Text_map.commit()
    rows = Text_map.get_dict_rows()
    rows

    Text_map.db_clear_table()
    rows = Text_map.get_dict_rows()

    Text_map.db_sync(['man in outside','man is rich 233', "man is rich"])
    Text_map.overide_map({"goodman": "great person"})
    Text_map.get_map()
    Text_map.delete_specific("attribute_value", [' guy is body'])
    Text_map.close()
